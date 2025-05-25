import streamlit as st
import openai
import json as js
import os

from PIL import Image
import PyPDF2
import pandas as pd
import lasio
import docx
import pytesseract

st.set_page_config(page_title="R2 - AI Learning Bot", layout="wide")
print("RUNNING THE CORRECT FILE! 🎉")  # For confirmation/debug

MEMORY_FILE = os.path.join(os.getcwd(), 'r2_memory.json')
FILES_KEY = "uploaded_files_store"
openai_api_key = st.text_input("Enter your OpenAI API key:", type="password")

st.title("R2 - AI Learning Bot")

# ---- SYSTEM PROMPT reflecting R2's real purpose ----
system_prompt = (
    "You are R2, the world's leading, continuously learning upstream oil and gas drilling AI assistant. "
    "Your mission is to empower drilling teams, engineers, and managers to optimize operations, minimize risk, and maximize well performance—globally and regionally. "
    "You continuously integrate, learn from, and reason over technical manuals, field reports, best practices, real-time EDR/WITSML/LAS data, user feedback, and uploaded documents. "
    "\n\n"
    "### Memory and Learning\n"
    "You retain and recall all information, conversations, and files uploaded within this session. Always act as if you have perfect session memory, synthesizing and summarizing as needed. Proactively reference earlier discussions, data, and documents when helpful."
    "\n\n"
    "### How to Respond\n"
    "- **Always format answers in Markdown:** use headers, bullet points, numbered lists, bold for key insights, and tables if appropriate.\n"
    "- Summarize and organize complex information clearly, with executive summaries if asked.\n"
    "- Provide actionable, technically correct, and step-by-step reasoning for drilling, completions, hydraulics, torque & drag, geomechanics, and optimization.\n"
    "- When reviewing files, documents, or logs, **cite data, explain context, and extract relevant facts**.\n"
    "- Adapt detail and language to the user's level (engineer, manager, trainee, etc.)—be concise but thorough, and clarify assumptions."
    "\n\n"
    "### Scope of Expertise\n"
    "- Drilling and completions (all methods)\n"
    "- BHA design, torque & drag, jar placement\n"
    "- Real-time dysfunction detection, optimization, automation\n"
    "- Wellbore stability, geomechanics, regional formation analysis\n"
    "- Regulatory, HSE, ESG, and field best practices\n"
    "- Cross-platform, collaborative workflows"
    "\n\n"
    "### Principles\n"
    "- Proactively prevent drilling hazards; recommend best practices\n"
    "- Provide explainable, referenced answers and calculations\n"
    "- Always state if information is missing or assumptions are needed\n"
    "- If asked, provide tables, diagrams, formulas, and references"
    "\n\n"
    "### Your Style\n"
    "- Professional, trustworthy, and adaptive\n"
    "- Always cite uploaded documents, session chat, and user data as sources\n"
    "- Encourage continuous learning, feedback, and knowledge sharing"
)

# --- Sidebar File Upload Section ---
with st.sidebar:
    st.markdown("### 📎 Uploaded Files")
    uploaded_files = st.file_uploader(
        "Upload any file (PDF, Word, Excel, CSV, LAS, images, etc):",
        accept_multiple_files=True,
        key="sidebar_file_upload"
    )
    # Save and store uploaded files in session_state
    if FILES_KEY not in st.session_state:
        st.session_state[FILES_KEY] = {}
    # Store new uploads in session_state
    if uploaded_files:
        for uploaded_file in uploaded_files:
            st.session_state[FILES_KEY][uploaded_file.name] = uploaded_file.getvalue()
    # List all uploaded files with option to remove
    for fname in list(st.session_state[FILES_KEY].keys()):
        st.markdown(
            f'<span style="font-size:0.95em; color:#444;">📄 <b>{fname}</b></span>',
            unsafe_allow_html=True
        )
        if st.button(f"Remove {fname}", key=f"rm_{fname}"):
            del st.session_state[FILES_KEY][fname]
            st.rerun()

# --- Load previous chat history if exists ---
if "history" not in st.session_state:
    if os.path.exists(MEMORY_FILE):
        with open(MEMORY_FILE, 'r') as f:
            st.session_state.history = js.load(f)
    else:
        st.session_state.history = []

# --- Main Area: Show the chat history ---
def render_chat(history):
    for msg in history:
        role = "🧑‍💻 You" if msg["role"] == "user" else "🤖 R2"
        # Add extra spacing and custom style for R2
        if msg["role"] == "assistant":
            st.markdown(
                f"""<div style='margin-top: 32px; margin-bottom: 30px; font-size:1.07em; color: #1a237e; 
                    font-family:Segoe UI,Arial,sans-serif; background:#f7f7ff; padding:12px 16px 10px 18px; border-radius:12px;'>
                <b>{role}:</b></div>""",
                unsafe_allow_html=True
            )
            # Show R2 reply as Markdown for headings/bullets/bold
            st.markdown(msg['content'], unsafe_allow_html=True)
        else:
            st.markdown(
                f"""<div style='margin-top: 18px; margin-bottom: 18px; font-size:1em; color:#333;'>
                <b>{role}:</b> {msg['content']}
                </div>""",
                unsafe_allow_html=True
            )

render_chat(st.session_state.history)
st.markdown("<br><hr><br>", unsafe_allow_html=True)

# --- Chat input form ---
with st.form(key="chat_form", clear_on_submit=True):
    user_input = st.text_input("You:", key="user_input", value="")
    submitted = st.form_submit_button("Send")

# --- Extraction logic, on upload or access ---
def extract_file_content(file_bytes, filename):
    name = filename.lower()
    from io import BytesIO
    uploaded_file = BytesIO(file_bytes)
    try:
        if name.endswith('.pdf'):
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() or ""
            return text
        elif name.endswith(('.jpg', '.jpeg', '.png')):
            image = Image.open(uploaded_file)
            text = pytesseract.image_to_string(image)
            return text
        elif name.endswith('.txt'):
            return uploaded_file.read().decode('utf-8', errors='ignore')
        elif name.endswith('.csv'):
            df = pd.read_csv(uploaded_file)
            return df.to_string()
        elif name.endswith(('.xls', '.xlsx')):
            df = pd.read_excel(uploaded_file)
            return df.to_string()
        elif name.endswith('.docx'):
            doc = docx.Document(uploaded_file)
            return '\n'.join([p.text for p in doc.paragraphs])
        elif name.endswith('.pptx'):
            from pptx import Presentation
            ppt = Presentation(uploaded_file)
            text = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        elif name.endswith('.las'):
            las = lasio.read(uploaded_file)
            return las.to_json()
        elif name.endswith('.json'):
            data = js.load(uploaded_file)
            return js.dumps(data, indent=2)
        elif name.endswith('.xml'):
            return uploaded_file.read().decode('utf-8', errors='ignore')
        elif name.endswith('.html'):
            return uploaded_file.read().decode('utf-8', errors='ignore')
        else:
            return None
    except Exception as e:
        return f"Error reading file: {e}"

# --- When chat is submitted ---
if submitted and user_input and openai_api_key:
    st.session_state.history.append({"role": "user", "content": user_input})

    # Gather all extracted file contents (for context only, not shown as chat message)
    file_contexts = []
    for fname, filebytes in st.session_state[FILES_KEY].items():
        file_text = extract_file_content(filebytes, fname)
        if file_text:
            file_contexts.append(f"File: {fname}\n{file_text[:1500]}")  # Truncate if very long

    # Build full AI context: R2's mission, files as reference, full chat
    chat_context = [
        {"role": "system", "content": system_prompt}
    ]
    if file_contexts:
        chat_context.append({"role": "system", "content": "The following files have been uploaded for reference:\n" + "\n\n".join(file_contexts)})
    chat_context += [ {"role": m["role"], "content": m["content"]} for m in st.session_state.history ]

    try:
        client = openai.OpenAI(api_key=openai_api_key)
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=chat_context,
        )
        answer = response.choices[0].message.content
    except Exception as e:
        answer = f"Error: {str(e)}"

    st.session_state.history.append({"role": "assistant", "content": answer})

    with open(MEMORY_FILE, 'w') as f:
        js.dump(st.session_state.history, f)
    st.rerun()

# --- Memory clear button ---
if st.button("Clear Memory"):
    st.session_state.history = []
    if os.path.exists(MEMORY_FILE):
        os.remove(MEMORY_FILE)
    st.rerun()